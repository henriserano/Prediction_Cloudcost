"""Generate a 2-slide pitch deck for the Sia FinOps AI platform."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Brand palette (neutral, replaceable once dropped into the Sia master)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0A, 0x1F, 0x44)
ACCENT = RGBColor(0xE6, 0x00, 0x28)
LIGHT = RGBColor(0xF2, 0xF4, 0xF8)
GREY = RGBColor(0x55, 0x5F, 0x6D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BAND = [
    RGBColor(0x0A, 0x1F, 0x44),
    RGBColor(0x14, 0x3A, 0x6E),
    RGBColor(0x1E, 0x55, 0x99),
    RGBColor(0x36, 0x7A, 0xB8),
    RGBColor(0x5B, 0x9B, 0xD5),
    RGBColor(0x8F, 0xBC, 0xE6),
    RGBColor(0xB8, 0xD4, 0xEE),
]


def add_textbox(slide, left, top, width, height, text, *, size=12,
                bold=False, color=NAVY, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return tb


def add_rect(slide, left, top, width, height, *, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, subtitle):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), fill=NAVY)
    add_rect(slide, Inches(0), Inches(0.85), Inches(13.333), Inches(0.05), fill=ACCENT)
    add_textbox(slide, Inches(0.4), Inches(0.12), Inches(11), Inches(0.5),
                title, size=22, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.4), Inches(0.5), Inches(11), Inches(0.35),
                subtitle, size=11, color=RGBColor(0xC7, 0xD2, 0xE5))
    # Sia mark placeholder
    add_textbox(slide, Inches(12.2), Inches(0.28), Inches(1), Inches(0.35),
                "Sia", size=14, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


def add_footer(slide, page_no):
    add_textbox(slide, Inches(0.4), Inches(7.15), Inches(10), Inches(0.3),
                "Sia. FinOps AI Platform. Confidential.", size=9, color=GREY)
    add_textbox(slide, Inches(12.5), Inches(7.15), Inches(0.6), Inches(0.3),
                f"{page_no}", size=9, color=GREY, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Deck
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

# ===========================================================================
# Slide 1. Value chain
# ===========================================================================
s1 = prs.slides.add_slide(blank_layout)
add_header(
    s1,
    "Plateforme FinOps IA. Cartographie fonctionnelle",
    "Une couverture end-to-end de la chaine de valeur d'un projet IA cloud",
)

# Value-chain band
chain = [
    ("1. Ingestion", "OAuth GCP. STS AWS.\nBigQuery export.\nCSV / Excel."),
    ("2. Preparation", "STL. Stationnarite.\nDistribution. ACF.\nMissingness."),
    ("3. Anomalies", "Outliers.\nDrift detection.\nAlertes precoces."),
    ("4. Modelisation", "6 modeles + ensemble.\nWalk-forward CV.\nIC 80% / 95%."),
    ("5. Simulation", "Cout projet IA.\nReco archi.\nMatrice de risques."),
    ("6. Consommation", "Dashboard Next.js.\nAgent MCP.\nHistorique conv."),
    ("7. Gouvernance", "API key. CSRF.\nCORS strict.\nSecScan CI."),
]

band_top = Inches(1.15)
band_height = Inches(2.1)
margin_x = Inches(0.4)
gap = Inches(0.1)
total_w = Inches(13.333) - 2 * margin_x
cell_w = (total_w - gap * (len(chain) - 1)) / len(chain)

for i, (title, body) in enumerate(chain):
    left = margin_x + (cell_w + gap) * i
    add_rect(s1, left, band_top, cell_w, Inches(0.5), fill=BAND[i])
    add_textbox(s1, left + Inches(0.1), band_top + Inches(0.08),
                cell_w - Inches(0.2), Inches(0.4),
                title, size=12, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s1, left, band_top + Inches(0.5), cell_w,
             band_height - Inches(0.5), fill=LIGHT)
    add_textbox(s1, left + Inches(0.12), band_top + Inches(0.6),
                cell_w - Inches(0.24), band_height - Inches(0.65),
                body, size=9, color=NAVY)

# Persona benefits below the band
persona_top = Inches(3.55)
persona_h = Inches(2.65)
personas = [
    ("Pour le CFO",
     "Prevision de facture a 90 jours avec fourchette d'incertitude.\n"
     "Chiffrage a priori d'un projet IA avant engagement budgetaire.\n"
     "Detection precoce des derives de consommation.",
     ACCENT),
    ("Pour le DSI",
     "Vision multi-cloud unifiee AWS + GCP.\n"
     "Agent conversationnel qui democratise le FinOps.\n"
     "Securite et IaC pretes pour audit.",
     NAVY),
    ("Pour le FinOps Lead",
     "EDA industrialisee, STL, drift et anomalies en un clic.\n"
     "Benchmark automatique de 6 modeles.\n"
     "Extensible via catalogue d'outils MCP.",
     RGBColor(0x1E, 0x55, 0x99)),
]

p_w = (total_w - gap * 2) / 3
for i, (title, body, color) in enumerate(personas):
    left = margin_x + (p_w + gap) * i
    add_rect(s1, left, persona_top, p_w, Inches(0.45), fill=color)
    add_textbox(s1, left + Inches(0.15), persona_top + Inches(0.06),
                p_w - Inches(0.3), Inches(0.35),
                title, size=13, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s1, left, persona_top + Inches(0.45), p_w,
             persona_h - Inches(0.45), fill=WHITE, line=color)
    add_textbox(s1, left + Inches(0.2), persona_top + Inches(0.6),
                p_w - Inches(0.4), persona_h - Inches(0.7),
                body, size=11, color=NAVY)

# Punch line
add_rect(s1, margin_x, Inches(6.4), total_w, Inches(0.55), fill=NAVY)
add_textbox(s1, margin_x + Inches(0.3), Inches(6.45),
            total_w - Inches(0.6), Inches(0.45),
            "Un seul outil, trois personas, sept etapes de la chaine de valeur IA. "
            "AWS et GCP couverts nativement.",
            size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s1, 1)

# ===========================================================================
# Slide 2. Cost simulation example
# ===========================================================================
s2 = prs.slides.add_slide(blank_layout)
add_header(
    s2,
    "Simulation projet IA. Exemple chiffre",
    "Assistant RAG interne. 500 utilisateurs. Deploiement GCP.",
)

# Left column. Hypotheses
add_rect(s2, Inches(0.4), Inches(1.15), Inches(4.0), Inches(0.4),
         fill=NAVY)
add_textbox(s2, Inches(0.55), Inches(1.19), Inches(3.7), Inches(0.32),
            "Hypotheses de scenario", size=13, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

hyp = [
    ("Cas d'usage", "Assistant RAG documentaire interne"),
    ("Utilisateurs actifs", "500 / mois"),
    ("Requetes moyennes", "40 par utilisateur / mois"),
    ("LLM principal", "Claude Sonnet 4.6"),
    ("Contexte moyen", "6k tokens in / 800 tokens out"),
    ("Vector store", "Managed. 1.2M chunks"),
    ("Deploiement", "GCP Cloud Run + Vertex"),
]
row_top = Inches(1.6)
row_h = Inches(0.32)
for i, (k, v) in enumerate(hyp):
    y = row_top + row_h * i
    bg = LIGHT if i % 2 == 0 else WHITE
    add_rect(s2, Inches(0.4), y, Inches(4.0), row_h, fill=bg)
    add_textbox(s2, Inches(0.55), y + Inches(0.04), Inches(1.7), Inches(0.28),
                k, size=10, bold=True, color=NAVY)
    add_textbox(s2, Inches(2.25), y + Inches(0.04), Inches(2.1), Inches(0.28),
                v, size=10, color=GREY)

# Middle column. Cost breakdown
add_rect(s2, Inches(4.55), Inches(1.15), Inches(4.2), Inches(0.4),
         fill=NAVY)
add_textbox(s2, Inches(4.7), Inches(1.19), Inches(3.9), Inches(0.32),
            "Breakdown de cout mensuel", size=13, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

items = [
    ("LLM Claude Sonnet 4.6 (in/out)", 1420, 0.42),
    ("Embeddings + reindex", 180, 0.05),
    ("Vector store managed", 260, 0.08),
    ("Cloud Run + reseau", 210, 0.06),
    ("Monitoring, logs, secrets", 90, 0.03),
    ("Marge d'incertitude (IC95)", 1220, 0.36),
]
total = sum(v for _, v, _ in items)

bar_top = Inches(1.6)
bar_h = Inches(0.36)
for i, (label, val, share) in enumerate(items):
    y = bar_top + bar_h * i
    add_textbox(s2, Inches(4.6), y + Inches(0.04), Inches(2.4), Inches(0.28),
                label, size=10, color=NAVY, bold=True)
    bar_w = Inches(1.0) * (share / 0.42) * 1.5
    add_rect(s2, Inches(6.95), y + Inches(0.08), max(bar_w, Inches(0.05)),
             Inches(0.2), fill=BAND[i % len(BAND)])
    add_textbox(s2, Inches(8.15), y + Inches(0.04), Inches(0.55), Inches(0.28),
                f"{val} EUR", size=10, color=NAVY, align=PP_ALIGN.RIGHT)

# Total line
add_rect(s2, Inches(4.55), bar_top + bar_h * len(items) + Inches(0.05),
         Inches(4.2), Inches(0.45), fill=NAVY)
add_textbox(s2, Inches(4.7),
            bar_top + bar_h * len(items) + Inches(0.09),
            Inches(2.5), Inches(0.35),
            "Total mensuel estime", size=12, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s2, Inches(7.3),
            bar_top + bar_h * len(items) + Inches(0.09),
            Inches(1.35), Inches(0.35),
            f"{total} EUR", size=13, bold=True, color=WHITE,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Right column. Reco + risks
add_rect(s2, Inches(8.9), Inches(1.15), Inches(4.05), Inches(0.4),
         fill=ACCENT)
add_textbox(s2, Inches(9.05), Inches(1.19), Inches(3.75), Inches(0.32),
            "Recommandation d'architecture", size=13, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

reco = [
    "Cache semantique sur les 20% de requetes recurrentes. "
    "Economie estimee 22%.",
    "Bascule vers Claude Haiku pour les requetes courtes classifiees. "
    "Economie estimee 15%.",
    "Batch d'embeddings quotidien plutot que streaming. "
    "Economie estimee 4%.",
]
r_top = Inches(1.6)
for i, txt in enumerate(reco):
    y = r_top + Inches(0.8) * i
    add_rect(s2, Inches(8.9), y, Inches(0.06), Inches(0.7), fill=ACCENT)
    add_textbox(s2, Inches(9.05), y + Inches(0.04),
                Inches(3.85), Inches(0.7),
                txt, size=10, color=NAVY)

# Risk matrix
add_rect(s2, Inches(8.9), Inches(4.05), Inches(4.05), Inches(0.4),
         fill=NAVY)
add_textbox(s2, Inches(9.05), Inches(4.09), Inches(3.75), Inches(0.32),
            "Matrice de risques", size=13, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)

risks = [
    ("Derive tokens en context", "Elevee", ACCENT),
    ("Cout embeddings a la reindex", "Moyen", RGBColor(0xF2, 0x99, 0x00)),
    ("Latence Vertex pic d'usage", "Moyen", RGBColor(0xF2, 0x99, 0x00)),
    ("Fuite donnees via logs", "Faible", RGBColor(0x2E, 0xA4, 0x4F)),
]
for i, (label, lvl, col) in enumerate(risks):
    y = Inches(4.5) + Inches(0.42) * i
    add_rect(s2, Inches(8.9), y, Inches(4.05), Inches(0.36), fill=LIGHT)
    add_textbox(s2, Inches(9.05), y + Inches(0.05),
                Inches(2.55), Inches(0.28),
                label, size=10, color=NAVY, bold=True)
    add_rect(s2, Inches(11.7), y + Inches(0.06), Inches(1.15),
             Inches(0.24), fill=col)
    add_textbox(s2, Inches(11.7), y + Inches(0.08),
                Inches(1.15), Inches(0.24),
                lvl, size=10, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# Bottom takeaway
add_rect(s2, Inches(0.4), Inches(6.4), Inches(12.55), Inches(0.55), fill=NAVY)
add_textbox(s2, Inches(0.7), Inches(6.45), Inches(12), Inches(0.45),
            "Sortie de la simulation, 3380 EUR / mois. Apres optimisations "
            "recommandees, 2020 EUR / mois. Payback outil < 1 mois.",
            size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

add_footer(s2, 2)


# ===========================================================================
# Slide 3. Connecter et comprendre vos donnees cloud
# ===========================================================================
def add_capability_slide(prs, page_no, title, subtitle, blocks, story):
    """Render a 'what you can do' slide with 3 vertical capability blocks and a
    concrete usage story at the bottom.

    blocks: list of 3 dicts {header, tagline, features (list[str]),
                             endpoints (list[str])}.
    story:  dict {persona, situation, action, outcome}.
    """
    s = prs.slides.add_slide(blank_layout)
    add_header(s, title, subtitle)

    margin_x = Inches(0.4)
    gap = Inches(0.2)
    total_w = Inches(13.333) - 2 * margin_x
    col_w = (total_w - gap * 2) / 3
    block_top = Inches(1.15)
    block_h = Inches(4.7)

    for i, blk in enumerate(blocks):
        left = margin_x + (col_w + gap) * i
        # header band
        add_rect(s, left, block_top, col_w, Inches(0.55), fill=NAVY)
        add_textbox(s, left + Inches(0.15), block_top + Inches(0.05),
                    col_w - Inches(0.3), Inches(0.25),
                    blk["header"], size=13, bold=True, color=WHITE)
        add_textbox(s, left + Inches(0.15), block_top + Inches(0.28),
                    col_w - Inches(0.3), Inches(0.22),
                    blk["tagline"], size=9,
                    color=RGBColor(0xC7, 0xD2, 0xE5))

        # body
        body_top = block_top + Inches(0.55)
        body_h = block_h - Inches(0.55)
        add_rect(s, left, body_top, col_w, body_h, fill=WHITE, line=NAVY)

        # "Vous pouvez" bullets
        add_textbox(s, left + Inches(0.15), body_top + Inches(0.1),
                    col_w - Inches(0.3), Inches(0.3),
                    "Vous pouvez", size=10, bold=True, color=ACCENT)

        y = body_top + Inches(0.4)
        for feat in blk["features"]:
            add_rect(s, left + Inches(0.18), y + Inches(0.09),
                     Inches(0.08), Inches(0.08), fill=ACCENT)
            add_textbox(s, left + Inches(0.35), y,
                        col_w - Inches(0.5), Inches(0.5),
                        feat, size=10, color=NAVY)
            y += Inches(0.42)

        # Endpoints strip
        y += Inches(0.05)
        add_rect(s, left + Inches(0.15), y, col_w - Inches(0.3),
                 Inches(0.02), fill=LIGHT)
        y += Inches(0.1)
        add_textbox(s, left + Inches(0.15), y, col_w - Inches(0.3),
                    Inches(0.22), "Points d'entree techniques",
                    size=8, bold=True, color=GREY)
        y += Inches(0.22)
        for ep in blk["endpoints"]:
            add_textbox(s, left + Inches(0.18), y, col_w - Inches(0.3),
                        Inches(0.2), ep, size=8, color=GREY)
            y += Inches(0.19)

    # Bottom user story band
    story_top = Inches(5.95)
    story_h = Inches(1.0)
    add_rect(s, margin_x, story_top, total_w, story_h, fill=LIGHT,
             line=NAVY)
    add_rect(s, margin_x, story_top, Inches(0.15), story_h, fill=ACCENT)

    add_textbox(s, margin_x + Inches(0.3), story_top + Inches(0.08),
                Inches(2.5), Inches(0.3),
                f"En pratique. {story['persona']}", size=11, bold=True,
                color=ACCENT)
    add_textbox(s, margin_x + Inches(0.3), story_top + Inches(0.35),
                total_w - Inches(0.5), Inches(0.6),
                f"{story['situation']} {story['action']} "
                f"Resultat. {story['outcome']}",
                size=10, color=NAVY)

    add_footer(s, page_no)
    return s


# --- Slide 3 -----------------------------------------------------------------
add_capability_slide(
    prs,
    page_no=3,
    title="Connecter, explorer, diagnostiquer",
    subtitle="Ce que vous pouvez faire des la premiere semaine d'usage",
    blocks=[
        {
            "header": "1. Brancher vos sources",
            "tagline": "Zero credential stockee cote client. Multi-cloud natif.",
            "features": [
                "Connecter GCP en un clic via OAuth2, sans partage de cle "
                "de service.",
                "Brancher AWS avec vos cles IAM validees par STS avant "
                "acceptation.",
                "Ingerer directement le BigQuery Billing Export officiel.",
                "Uploader un CSV ou Excel pour un POC rapide sans "
                "integration IT.",
                "Isoler les credentials par utilisateur dans un coffre "
                "chiffre.",
            ],
            "endpoints": [
                "/api/gcp/auth. /api/gcp/callback. /api/gcp/projects.",
                "/api/aws/status. /api/aws/connect.",
                "/api/credentials. /api/data/status.",
            ],
        },
        {
            "header": "2. Explorer vos donnees",
            "tagline": "Une EDA industrialisee, jusqu'ici reservee aux "
                       "data scientists.",
            "features": [
                "Visualiser vos KPI FinOps consolides sur une periode "
                "arbitraire.",
                "Decomposer votre serie de cout en tendance, "
                "saisonnalite et residu (STL).",
                "Mesurer la stationnarite (test ADF) avant toute "
                "prevision.",
                "Cartographier la repartition par service et sous-service "
                "GCP ou AWS.",
                "Detecter les valeurs manquantes et les incoherences "
                "de scaling.",
            ],
            "endpoints": [
                "/api/analytics/kpi. /api/analytics/daily.",
                "/api/analytics/stl. /api/analytics/stationarity.",
                "/api/analytics/services. /api/analytics/missing.",
            ],
        },
        {
            "header": "3. Diagnostiquer",
            "tagline": "Passer de l'alerte binaire a la comprehension "
                       "statistique.",
            "features": [
                "Reperer les anomalies temporelles avec seuil "
                "statistiquement justifie.",
                "Detecter les outliers multivaries sur plusieurs "
                "services simultanes.",
                "Identifier les drifts distributionnels avant qu'ils "
                "explosent la facture.",
                "Comparer la distribution de cout entre deux periodes "
                "ou deux equipes.",
                "Consulter les Cloud Logs GCP directement dans "
                "l'interface.",
            ],
            "endpoints": [
                "/api/analytics/anomalies. /api/analytics/outliers.",
                "/api/analytics/drift. /api/analytics/distribution.",
                "/api/gcp/logs. /api/gcp/services.",
            ],
        },
    ],
    story={
        "persona": "Un analyste FinOps chez un retailer, semaine 1.",
        "situation": "Il branche GCP en 2 minutes et importe l'historique "
                     "billing sur 18 mois.",
        "action": "En moins d'une heure il visualise la decomposition STL "
                  "qui revele une saisonnalite hebdo jusqu'ici masquee "
                  "par la moyenne mensuelle, et repere deux drifts "
                  "de service concomitants du dernier deploiement produit.",
        "outcome": "3 heures de travail evitent une derive mensuelle "
                   "estimee a 8 000 EUR.",
    },
)


# --- Slide 4 -----------------------------------------------------------------
add_capability_slide(
    prs,
    page_no=4,
    title="Anticiper, comparer, arbitrer",
    subtitle="Ce que vous pouvez faire quand vous devez decider et "
             "planifier",
    blocks=[
        {
            "header": "4. Prevoir la facture",
            "tagline": "Une prevision honnete avec fourchette d'incertitude "
                       "explicite.",
            "features": [
                "Generer un forecast a 7, 30, 60 ou 90 jours sur "
                "n'importe quelle serie.",
                "Comparer six modeles statistiques evalues en "
                "walk-forward cross-validation.",
                "Choisir manuellement un modele ou laisser un ensemble "
                "pondere decider.",
                "Consulter les intervalles de confiance 80% et 95% "
                "sur chaque point.",
                "Exporter les previsions pour votre outil de budget.",
            ],
            "endpoints": [
                "/api/forecast. /api/forecast/summary.",
                "/api/forecast/models. /api/analytics/ensemble-forecast.",
                "Modeles: ETS, Theta, ARIMA, SES, Holt-Winters, SNaive.",
            ],
        },
        {
            "header": "5. Simuler un projet IA",
            "tagline": "Chiffrer avant d'engager. Comparer les "
                       "architectures.",
            "features": [
                "Estimer le cout mensuel d'un projet IA a partir de "
                "vos hypotheses metier.",
                "Choisir votre LLM dans un catalogue de pricing "
                "actualise.",
                "Modeliser tools, RAG, agents, deploiement, volume "
                "d'utilisateurs.",
                "Recevoir 2 a 3 recommandations d'architecture avec "
                "gain chiffre.",
                "Obtenir une matrice de risques et les axes "
                "d'optimisation prioritaires.",
            ],
            "endpoints": [
                "/api/simulate. /api/simulate/reference-catalog.",
                "Modules: CostBreakdown, ArchitectureRecommendation.",
                "RiskAssessment. ProjectedEvents.",
            ],
        },
        {
            "header": "6. Arbitrer les scenarios",
            "tagline": "Un tableau de bord pour trancher entre plusieurs "
                       "options.",
            "features": [
                "Comparer plusieurs scenarios simules cote a cote.",
                "Sauvegarder les scenarios pour les rejouer avec "
                "de nouvelles hypotheses.",
                "Visualiser l'impact d'une reduction de tokens, d'un "
                "changement de modele.",
                "Exporter le comparatif en support de comite "
                "d'investissement.",
                "Tracer les evenements coup projetes sur les 12 "
                "prochains mois.",
            ],
            "endpoints": [
                "/api/simulate. /simulation UI.",
                "/api/analytics/scaling. /api/analytics/dim-reduction.",
                "Historique par session utilisateur.",
            ],
        },
    ],
    story={
        "persona": "Une DSI d'assureur qui hesite entre trois "
                   "architectures RAG.",
        "situation": "Elle doit trancher en comite d'investissement dans "
                     "10 jours.",
        "action": "Elle lance trois simulations en 15 minutes, compare "
                  "le cout mensuel, les risques et les recommandations "
                  "d'archi. La solution isole une reco de cache semantique "
                  "qui divise par deux le poste LLM sans degrader la "
                  "qualite.",
        "outcome": "Decision documentee, chiffree, defendable devant "
                   "COMEX. Payback outil sous 3 semaines.",
    },
)


# --- Slide 5 -----------------------------------------------------------------
add_capability_slide(
    prs,
    page_no=5,
    title="Piloter au quotidien via un agent IA",
    subtitle="Ce que vous pouvez faire quand vous voulez industrialiser "
             "l'usage FinOps",
    blocks=[
        {
            "header": "7. Assistant conversationnel",
            "tagline": "Un copilote FinOps en langage naturel, ouvert "
                       "aux integrations.",
            "features": [
                "Poser vos questions FinOps en francais ou anglais, "
                "sans SQL ni BigQuery.",
                "Laisser l'agent invoquer les outils analytiques, "
                "forecast, simulation.",
                "Historiser les conversations par utilisateur, "
                "reprendre le fil d'analyse.",
                "Exposer les outils au standard MCP pour integrer "
                "ChatGPT Enterprise, Claude, Copilot.",
                "Enchainer les analyses. L'agent appelle STL puis "
                "forecast puis simule.",
            ],
            "endpoints": [
                "/api/chat. /api/conversations.",
                "/api/tools. /api/tools/invoke. Standard MCP.",
                "Compatible LLM Enterprise du client.",
            ],
        },
        {
            "header": "8. Gouvernance et securite",
            "tagline": "Prete pour l'audit d'un DSI ou d'un RSSI.",
            "features": [
                "Authentifier vos utilisateurs par PIN et session "
                "cookie httpOnly.",
                "Proteger les endpoints critiques par cle d'API "
                "dediee (X-API-Key).",
                "Cloisonner CORS strictement en production, refus "
                "de demarrer si mal configure.",
                "Journaliser les evenements applicatifs pour audit "
                "et facturation interne.",
                "Ne jamais stocker de tokens OAuth partages. Chaque "
                "session est isolee.",
            ],
            "endpoints": [
                "/api/auth. /api/events. SEC-013, SEC-014, SEC-015.",
                "CI: pip-audit, npm audit, bandit, semgrep.",
                "Cache admin purgeable par cle.",
            ],
        },
        {
            "header": "9. Deploiement et TCO",
            "tagline": "De la demo au run production sans reecrire.",
            "features": [
                "Deployer en un script sur AWS (ECS Fargate + ALB + "
                "ECR + CloudWatch).",
                "Beneficier d'une IaC Terraform auditable, un fichier "
                "par concern.",
                "Monter en charge horizontalement selon vos volumes "
                "de facture.",
                "Integrer la CI/CD du client, badges de securite "
                "sur chaque commit.",
                "Personnaliser la marque et le catalogue de modeles "
                "selon vos standards.",
            ],
            "endpoints": [
                "deploy.sh. terraform apply.",
                ".github/workflows/*. Docker smoke tests.",
                "Marque Sia. Extensible par accompagnement.",
            ],
        },
    ],
    story={
        "persona": "Un CFO adjoint qui reprend le pilotage FinOps groupe.",
        "situation": "Il herite d'un perimetre AWS + GCP eclate entre "
                     "5 BU, aucun langage commun.",
        "action": "Ses equipes deploient la plateforme en 2 semaines. "
                  "Chaque responsable BU utilise l'agent pour poser ses "
                  "questions, la simulation pour arbitrer les nouveaux "
                  "projets IA, et le forecast pour son budget annuel. "
                  "La gouvernance passe l'audit RSSI sans reserve.",
        "outcome": "Vision consolidee groupe en < 30 jours. Reduction "
                   "de 12 a 18 % de la facture cloud sur 6 mois.",
    },
)


# ===========================================================================
# Slide 6. Cadrage / Build / Run
# ===========================================================================
s6 = prs.slides.add_slide(blank_layout)
add_header(
    s6,
    "Impact sur le cycle de vie d'un projet IA",
    "Positionnement de la plateforme sur les phases Cadrage, Build et Run",
)

margin_x = Inches(0.4)
gap = Inches(0.2)
total_w = Inches(13.333) - 2 * margin_x

# Timeline arrow visual at the top
tl_top = Inches(1.15)
tl_h = Inches(0.55)
add_rect(s6, margin_x, tl_top, total_w, tl_h, fill=LIGHT)

phases = [
    ("CADRAGE", "Chiffrer et arbitrer avant d'engager le budget",
     BAND[0]),
    ("BUILD", "Piloter les couts au fil du developpement",
     BAND[2]),
    ("RUN", "Anticiper et proteger le budget en production",
     BAND[4]),
]

seg_w = total_w / 3
for i, (ph, tag, col) in enumerate(phases):
    left = margin_x + seg_w * i
    add_rect(s6, left + Inches(0.05), tl_top + Inches(0.05),
             seg_w - Inches(0.1), tl_h - Inches(0.1), fill=col)
    add_textbox(s6, left + Inches(0.2), tl_top + Inches(0.05),
                Inches(1.3), Inches(0.25),
                ph, size=13, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s6, left + Inches(1.5), tl_top + Inches(0.05),
                seg_w - Inches(1.6), tl_h - Inches(0.1),
                tag, size=10, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)

# Content blocks for each phase
block_top = Inches(1.85)
block_h = Inches(4.1)
col_w = (total_w - gap * 2) / 3

phases_detail = [
    {
        "color": BAND[0],
        "questions_title": "Les questions auxquelles vous repondez",
        "questions": [
            "Combien ce projet IA va-t-il vraiment couter par mois ?",
            "Quelle architecture retenir. Cloud Run, VM, "
            "serverless ?",
            "Quel LLM. Sonnet, Haiku, open source auto-heberge ?",
            "Quels sont les risques a documenter pour le COMEX ?",
        ],
        "tools_title": "Ce que la plateforme apporte",
        "tools": [
            "Simulateur de cout projet IA (LLM, tools, users, "
            "deploiement).",
            "Catalogue de pricing LLM et outils actualise.",
            "Recommandations d'architecture chiffrees.",
            "Matrice de risques automatique.",
            "Export du business case pour comite "
            "d'investissement.",
        ],
        "outcome_title": "Livrable pour l'equipe",
        "outcome": "Un business case chiffre, defendable et versionne "
                   "en moins d'une journee, la ou il fallait 2 a "
                   "3 semaines de conseil.",
        "kpi": "-70% de temps de cadrage",
    },
    {
        "color": BAND[2],
        "questions_title": "Les questions auxquelles vous repondez",
        "questions": [
            "Sommes-nous dans l'enveloppe validee au cadrage ?",
            "Ou partent les tokens sur cette semaine de sprint ?",
            "Le cout unitaire par requete s'ameliore ou se degrade ?",
            "Une derive s'est-elle installee entre deux commits ?",
        ],
        "tools_title": "Ce que la plateforme apporte",
        "tools": [
            "Suivi journalier des couts dev, staging et prod par "
            "service.",
            "Detection precoce de drift entre deux periodes de "
            "sprint.",
            "Anomalies et outliers detectes des l'apparition.",
            "Reforecast automatique a chaque nouvelle donnee.",
            "Agent conversationnel pour investiguer une derive "
            "sans SQL.",
        ],
        "outcome_title": "Livrable pour l'equipe",
        "outcome": "Un point FinOps hebdomadaire objectif, integre au "
                   "rituel d'equipe. Les optimisations LLM sont "
                   "prioritisees en connaissance de cause, pas au "
                   "doigt mouille.",
        "kpi": "-25% de derive moyenne",
    },
    {
        "color": BAND[4],
        "questions_title": "Les questions auxquelles vous repondez",
        "questions": [
            "Combien va couter le projet le mois prochain ?",
            "Une anomalie de facturation est-elle en cours ?",
            "Une nouvelle campagne va-t-elle exploser mon budget ?",
            "Ou optimiser en priorite sans degrader la qualite ?",
        ],
        "tools_title": "Ce que la plateforme apporte",
        "tools": [
            "Forecast 30 / 60 / 90 jours avec intervalles "
            "de confiance.",
            "Alertes anomalies et drift, avant l'arrivee de la "
            "facture.",
            "Comparaison realise vs prevu, mois par mois.",
            "Historique des conversations d'investigation pour "
            "audit.",
            "Journalisation d'evenements exportable vers "
            "l'outillage FinOps existant.",
        ],
        "outcome_title": "Livrable pour l'equipe",
        "outcome": "Un pilotage continu et documente, avec un delai "
                   "d'alerte de plusieurs jours voire semaines "
                   "avant l'incident financier.",
        "kpi": "-12 a -18% sur la facture 6 mois",
    },
]

for i, ph in enumerate(phases_detail):
    left = margin_x + (col_w + gap) * i
    # Left color band
    add_rect(s6, left, block_top, Inches(0.12), block_h, fill=ph["color"])
    # Card background
    add_rect(s6, left + Inches(0.12), block_top,
             col_w - Inches(0.12), block_h, fill=WHITE, line=NAVY)

    # Section 1: Questions
    y = block_top + Inches(0.15)
    add_textbox(s6, left + Inches(0.28), y,
                col_w - Inches(0.4), Inches(0.25),
                ph["questions_title"], size=10, bold=True,
                color=ACCENT)
    y += Inches(0.3)
    for q in ph["questions"]:
        add_textbox(s6, left + Inches(0.28), y,
                    col_w - Inches(0.4), Inches(0.4),
                    f"? {q}", size=9, color=NAVY)
        y += Inches(0.32)

    # Divider
    y += Inches(0.02)
    add_rect(s6, left + Inches(0.28), y, col_w - Inches(0.4),
             Inches(0.02), fill=LIGHT)
    y += Inches(0.1)

    # Section 2: Tools
    add_textbox(s6, left + Inches(0.28), y,
                col_w - Inches(0.4), Inches(0.25),
                ph["tools_title"], size=10, bold=True, color=NAVY)
    y += Inches(0.28)
    for t in ph["tools"]:
        add_rect(s6, left + Inches(0.3), y + Inches(0.08),
                 Inches(0.07), Inches(0.07), fill=ph["color"])
        add_textbox(s6, left + Inches(0.45), y,
                    col_w - Inches(0.55), Inches(0.4),
                    t, size=9, color=NAVY)
        y += Inches(0.3)

# Bottom impact band with 3 quantified KPIs
band_top = Inches(6.1)
band_h = Inches(0.85)
add_rect(s6, margin_x, band_top, total_w, band_h, fill=NAVY)
add_textbox(s6, margin_x + Inches(0.25), band_top + Inches(0.08),
            Inches(3), Inches(0.3),
            "Impact quantifie", size=11, bold=True, color=ACCENT)

impact_cell_w = (total_w - Inches(0.6)) / 3
for i, ph in enumerate(phases_detail):
    left = margin_x + Inches(0.25) + impact_cell_w * i
    add_textbox(s6, left, band_top + Inches(0.35),
                impact_cell_w - Inches(0.2), Inches(0.28),
                phases[i][0], size=10, bold=True,
                color=RGBColor(0xC7, 0xD2, 0xE5))
    add_textbox(s6, left, band_top + Inches(0.55),
                impact_cell_w - Inches(0.2), Inches(0.28),
                ph["kpi"], size=14, bold=True, color=WHITE)

add_footer(s6, 6)


# ---------------------------------------------------------------------------
out = "/Users/henri/Vscode-DRAFT/Facturation_prediction/Sia_FinOps_AI_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
